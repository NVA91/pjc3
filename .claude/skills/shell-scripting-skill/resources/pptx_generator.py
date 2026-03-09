#!/usr/bin/env python3
# LADESTUFE 3 — RESSOURCE | KATEGORIE: Hybride Workflows / Office-Arbeit
# Quellcode belastet das Kontextfenster nicht. Nur stdout (Ergebnis) wird
# dem Modell zurückgegeben. Ausführung in Sandbox auf expliziten Nutzerruf.
#
# BESCHREIBUNG: PowerPoint-Generator für hybride Office/Programmier-Workflows.
#               Erstellt präsentationsfertige .pptx-Dateien aus Excel-Workbooks,
#               SVG/PNG-Diagrammen, CSV-Daten und Markdown-Text.
#               Unterstützt Corporate-Templates und automatische KPI-Layouts.
#
# VERWENDUNG:
#   python pptx_generator.py erstellen TITEL [OPTIONEN]
#   python pptx_generator.py hinzufuegen PFAD SLIDE-TYP [OPTIONEN]
#   python pptx_generator.py vorschau   PFAD
#   python pptx_generator.py template   PFAD --output PFAD
#   python pptx_generator.py konvertieren PFAD --format pdf|html
#
# SLIDE-TYPEN:
#   titel      — Titelfolie (Titel + Untertitel + Datum)
#   kpi        — KPI-Übersicht (bis zu 6 Kennzahlen mit Trend-Pfeil)
#   tabelle    — Datentabelle aus CSV oder Excel-Sheet
#   diagramm   — SVG/PNG-Diagramm vollflächig oder mit Beschriftung
#   text       — Freitext / Stichpunkte aus Markdown
#   abschnitt  — Abschnittstrenner (farbiger Hintergrund + Kapitelname)
#   leer       — Leere Folie mit optionalem Hintergrundgradient
#
# VORAUSSETZUNG:
#   pip install python-pptx pillow cairosvg
#   cairosvg nur für SVG→PNG-Konvertierung; ohne cairosvg: nur PNG-Einbettung
#
# HYBRID-WORKFLOW (Quartalsbericht-Pipeline):
#   python pptx_generator.py erstellen "Q1 2026 Quartalsbericht" \
#     --excel q1_bericht.xlsx \
#     --diagramm transaktionsfluss.svg \
#     --output Q1_2026_Quartalsbericht.pptx

from __future__ import annotations

import sys
import os
import csv
import json
import io
from pathlib import Path
from datetime import datetime
from typing import Any

# =============================================================================
# KONFIGURATION
# =============================================================================

DEFAULT_BREITE  = 10    # Zoll (Standard-Widescreen 16:9)
DEFAULT_HOEHE   = 5.625 # Zoll
DEFAULT_THEMA   = "dunkel"  # dunkel | hell | corporate

# Farbpalette (angepasst an Design-Manifest: kein generisches Office-Blau)
THEMEN: dict[str, dict[str, str]] = {
    "dunkel": {
        "hintergrund":   "0A0A0F",
        "hintergrund2":  "12121A",
        "akzent":        "7C3AED",
        "akzent2":       "F59E0B",
        "text":          "FAFAF9",
        "text_muted":    "71717A",
        "border":        "27272A",
        "positiv":       "22C55E",
        "negativ":       "EF4444",
        "neutral":       "3B82F6",
    },
    "hell": {
        "hintergrund":   "FAFAFA",
        "hintergrund2":  "F4F4F5",
        "akzent":        "7C3AED",
        "akzent2":       "F59E0B",
        "text":          "09090B",
        "text_muted":    "71717A",
        "border":        "E4E4E7",
        "positiv":       "16A34A",
        "negativ":       "DC2626",
        "neutral":       "2563EB",
    },
    "corporate": {
        "hintergrund":   "1E293B",
        "hintergrund2":  "0F172A",
        "akzent":        "3B82F6",
        "akzent2":       "10B981",
        "text":          "F8FAFC",
        "text_muted":    "94A3B8",
        "border":        "334155",
        "positiv":       "10B981",
        "negativ":       "EF4444",
        "neutral":       "3B82F6",
    },
}

# =============================================================================
# UTILITIES
# =============================================================================

def _ts() -> str:
    return datetime.utcnow().strftime("%Y-%m-%dT%H:%M:%SZ")

def _log(msg: str) -> None:
    print(f"[{_ts()}] {msg}", file=sys.stderr)

def _ok(msg: str) -> None:
    print(f"OK  | {msg}", file=sys.stderr)

def _err(msg: str, code: int = 1) -> None:
    print(f"ERR | {msg}", file=sys.stderr)
    sys.exit(code)

def _json_out(daten: Any) -> None:
    print(json.dumps(daten, ensure_ascii=False, indent=2, default=str))

def _pptx_import():
    """Importiert python-pptx mit verständlicher Fehlermeldung."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt
        return Presentation, Inches, Pt, Emu, RGBColor, PP_ALIGN
    except ImportError:
        _err(
            "python-pptx nicht installiert.\n"
            "Installation: pip install python-pptx\n"
            "Für SVG-Support zusätzlich: pip install cairosvg pillow"
        )

def _hex_zu_rgb(hex_str: str):
    """Konvertiert Hex-Farbe in RGBColor."""
    _, _, _, RGBColor, _ = _pptx_import()[1:]
    hex_str = hex_str.lstrip("#")
    return RGBColor(int(hex_str[0:2], 16), int(hex_str[2:4], 16), int(hex_str[4:6], 16))

def _svg_zu_png(svg_pfad: str) -> bytes | None:
    """Konvertiert SVG in PNG-Bytes. Gibt None zurück falls cairosvg fehlt."""
    try:
        import cairosvg
        return cairosvg.svg2png(url=svg_pfad, output_width=1920, output_height=1080)
    except ImportError:
        _log("cairosvg nicht installiert — SVG wird als Pfad referenziert statt eingebettet")
        _log("Installation: pip install cairosvg (benötigt cairo-Bibliothek)")
        return None

def _csv_lesen(pfad: str, max_zeilen: int = 20) -> tuple[list[str], list[list[str]]]:
    """Liest CSV und gibt (Header, Zeilen[:max_zeilen]) zurück."""
    with open(pfad, encoding="utf-8-sig", newline="") as f:
        reader = csv.reader(f)
        header = next(reader, [])
        zeilen = [row for i, row in enumerate(reader) if i < max_zeilen]
    return header, zeilen

def _excel_sheet_lesen(excel_pfad: str, sheet: str | None = None,
                        max_zeilen: int = 20) -> tuple[list[str], list[list]]:
    """Liest ein Excel-Sheet (openpyxl)."""
    try:
        import openpyxl
    except ImportError:
        _err("openpyxl nicht installiert: pip install openpyxl")
    wb = openpyxl.load_workbook(excel_pfad, read_only=True, data_only=True)
    ws = wb[sheet] if sheet else wb.active
    zeilen = list(ws.iter_rows(values_only=True))
    wb.close()
    if not zeilen:
        return [], []
    header = [str(v) if v is not None else "" for v in zeilen[0]]
    daten  = [[str(v) if v is not None else "" for v in row]
              for row in zeilen[1:max_zeilen+1]]
    return header, daten

# =============================================================================
# PRÄSENTATION — Kernfunktionen
# =============================================================================

class Praesentation:
    """Wrapper um python-pptx Presentation mit Design-Manifest-Stilen."""

    def __init__(self, thema: str = DEFAULT_THEMA, template: str | None = None):
        Presentation, Inches, Pt, Emu, RGBColor, PP_ALIGN = _pptx_import()
        self.prs = Presentation(template) if template else Presentation()
        self.prs.slide_width  = Inches(DEFAULT_BREITE)
        self.prs.slide_height = Inches(DEFAULT_HOEHE)
        self.farben = THEMEN.get(thema, THEMEN["dunkel"])
        self._Inches = Inches
        self._Pt = Pt
        self._RGBColor = RGBColor
        self._PP_ALIGN = PP_ALIGN
        self.slide_count = 0

    def _layout(self, idx: int = 6):
        """Leeres Layout (blank)."""
        return self.prs.slide_layouts[min(idx, len(self.prs.slide_layouts)-1)]

    def _hintergrund_setzen(self, slide, farbe_hex: str):
        """Setzt Folienhintergrund auf Vollfarbe."""
        from pptx.util import Pt
        from pptx.dml.color import RGBColor
        from pptx.oxml.ns import qn
        from lxml import etree
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(*[int(farbe_hex[i:i+2], 16) for i in (0,2,4)])

    def _textbox(self, slide, text: str, x: float, y: float,
                  breite: float, hoehe: float,
                  schriftgroesse: int = 18,
                  fett: bool = False,
                  kursiv: bool = False,
                  farbe_hex: str | None = None,
                  ausrichtung: str = "links",
                  zeilenabstand: float | None = None):
        """Fügt Textbox auf Folie ein."""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        Inches_ = self._Inches

        txb = slide.shapes.add_textbox(
            Inches_(x), Inches_(y), Inches_(breite), Inches_(hoehe)
        )
        tf = txb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        if ausrichtung == "mitte":
            p.alignment = PP_ALIGN.CENTER
        elif ausrichtung == "rechts":
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT

        run = p.add_run()
        run.text = text
        run.font.size = Pt(schriftgroesse)
        run.font.bold = fett
        run.font.italic = kursiv

        farbe = farbe_hex or self.farben["text"]
        run.font.color.rgb = RGBColor(*[int(farbe[i:i+2], 16) for i in (0,2,4)])

        return txb

    def _rechteck(self, slide, x: float, y: float, breite: float, hoehe: float,
                   farbe_hex: str, ecke_radius: float = 0):
        """Fügt farbiges Rechteck ein."""
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            self._Inches(x), self._Inches(y),
            self._Inches(breite), self._Inches(hoehe)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(*[int(farbe_hex[i:i+2], 16) for i in (0,2,4)])
        shape.line.fill.background()
        return shape

    def slide_titel(self, titel: str, untertitel: str = "",
                    datum: str | None = None, slide_nummer: bool = True):
        """Titelfolie mit großem Titel, Untertitel und optionalem Datum."""
        slide = self.prs.slides.add_slide(self._layout())
        self._hintergrund_setzen(slide, self.farben["hintergrund"])

        # Akzentbalken oben
        self._rechteck(slide, 0, 0, 10, 0.05, self.farben["akzent"])

        # Titel
        self._textbox(slide, titel,
                       x=0.6, y=1.8, breite=8.8, hoehe=1.5,
                       schriftgroesse=44, fett=True,
                       farbe_hex=self.farben["text"])

        # Untertitel
        if untertitel:
            self._textbox(slide, untertitel,
                           x=0.6, y=3.4, breite=8.8, hoehe=0.8,
                           schriftgroesse=20,
                           farbe_hex=self.farben["text_muted"])

        # Datum
        if datum:
            self._textbox(slide, datum,
                           x=0.6, y=4.8, breite=4, hoehe=0.4,
                           schriftgroesse=12,
                           farbe_hex=self.farben["akzent2"])

        self.slide_count += 1
        _ok(f"Folie {self.slide_count}: Titelfolie '{titel}'")

    def slide_abschnitt(self, kapitel: str, nummer: str = ""):
        """Abschnittstrenner mit farbigem Hintergrund."""
        slide = self.prs.slides.add_slide(self._layout())
        self._hintergrund_setzen(slide, self.farben["akzent"])

        if nummer:
            self._textbox(slide, nummer,
                           x=0.6, y=1.5, breite=9, hoehe=0.8,
                           schriftgroesse=60, fett=True,
                           farbe_hex=self.farben["hintergrund"] + "80")

        self._textbox(slide, kapitel,
                       x=0.6, y=2.0, breite=9, hoehe=1.5,
                       schriftgroesse=36, fett=True,
                       farbe_hex=self.farben["text"],
                       ausrichtung="mitte")

        self.slide_count += 1
        _ok(f"Folie {self.slide_count}: Abschnitt '{kapitel}'")

    def slide_kpi(self, kpis: list[dict]):
        """
        KPI-Übersichtsfolie mit bis zu 6 Kennzahlen.
        kpis = [{"label": "Umsatz", "wert": "€ 2.4M", "trend": "+12%", "positiv": True}, ...]
        """
        slide = self.prs.slides.add_slide(self._layout())
        self._hintergrund_setzen(slide, self.farben["hintergrund"])
        self._rechteck(slide, 0, 0, 10, 0.04, self.farben["akzent"])

        kpis = kpis[:6]
        n = len(kpis)
        spalten = 3 if n > 3 else n
        zeilen  = 2 if n > 3 else 1

        kachel_b = 2.8 if spalten == 3 else 3.8
        kachel_h = 1.8 if zeilen == 2 else 2.8
        start_x  = (10 - spalten * (kachel_b + 0.2)) / 2
        start_y  = 0.8

        for i, kpi in enumerate(kpis):
            col  = i % spalten
            row  = i // spalten
            x    = start_x + col * (kachel_b + 0.25)
            y    = start_y + row * (kachel_h + 0.2)

            # Kachel-Hintergrund
            self._rechteck(slide, x, y, kachel_b, kachel_h, self.farben["hintergrund2"])

            # Label
            self._textbox(slide, kpi.get("label", ""),
                           x=x+0.15, y=y+0.15, breite=kachel_b-0.3, hoehe=0.4,
                           schriftgroesse=11,
                           farbe_hex=self.farben["text_muted"])

            # Wert
            self._textbox(slide, kpi.get("wert", ""),
                           x=x+0.15, y=y+0.5, breite=kachel_b-0.3, hoehe=0.7,
                           schriftgroesse=28, fett=True,
                           farbe_hex=self.farben["text"])

            # Trend
            trend = kpi.get("trend", "")
            if trend:
                positiv = kpi.get("positiv", True)
                trend_farbe = self.farben["positiv"] if positiv else self.farben["negativ"]
                self._textbox(slide, trend,
                               x=x+0.15, y=y+1.2, breite=kachel_b-0.3, hoehe=0.4,
                               schriftgroesse=13, fett=True,
                               farbe_hex=trend_farbe)

        self.slide_count += 1
        _ok(f"Folie {self.slide_count}: KPI-Übersicht ({len(kpis)} Kennzahlen)")

    def slide_tabelle(self, header: list[str], zeilen: list[list[str]],
                       titel: str = "", max_zeilen: int = 15):
        """Datentabelle aus Header + Zeilen."""
        slide = self.prs.slides.add_slide(self._layout())
        self._hintergrund_setzen(slide, self.farben["hintergrund"])
        self._rechteck(slide, 0, 0, 10, 0.04, self.farben["akzent"])

        if titel:
            self._textbox(slide, titel,
                           x=0.3, y=0.1, breite=9.4, hoehe=0.5,
                           schriftgroesse=18, fett=True,
                           farbe_hex=self.farben["text"])

        zeilen = zeilen[:max_zeilen]
        n_cols  = len(header)
        n_rows  = len(zeilen) + 1  # +1 für Header
        start_y = 0.7 if titel else 0.3
        t_hoehe = 5.625 - start_y - 0.2
        t_breite= 9.4

        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor

        tbl = slide.shapes.add_table(
            n_rows, n_cols,
            self._Inches(0.3), self._Inches(start_y),
            self._Inches(t_breite), self._Inches(t_hoehe)
        ).table

        # Header-Zeile
        for j, h in enumerate(header):
            cell = tbl.cell(0, j)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(*[int(self.farben["akzent"][i:i+2], 16) for i in (0,2,4)])
            p = cell.text_frame.paragraphs[0]
            p.runs[0].font.bold = True
            p.runs[0].font.size = Pt(10)
            p.runs[0].font.color.rgb = RGBColor(250, 250, 249)

        # Daten-Zeilen
        for i, zeile in enumerate(zeilen):
            bg = self.farben["hintergrund"] if i % 2 == 0 else self.farben["hintergrund2"]
            for j, val in enumerate(zeile[:n_cols]):
                cell = tbl.cell(i+1, j)
                cell.text = str(val)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(*[int(bg[i2:i2+2], 16) for i2 in (0,2,4)])
                p = cell.text_frame.paragraphs[0]
                if p.runs:
                    p.runs[0].font.size = Pt(9)
                    p.runs[0].font.color.rgb = RGBColor(*[int(self.farben["text"][i2:i2+2], 16) for i2 in (0,2,4)])

        self.slide_count += 1
        _ok(f"Folie {self.slide_count}: Tabelle '{titel}' ({len(zeilen)} Zeilen, {n_cols} Spalten)")

    def slide_diagramm(self, bild_pfad: str, titel: str = "",
                        vollflächig: bool = True, beschriftung: str = ""):
        """Diagramm-Folie (SVG wird zu PNG konvertiert, PNG direkt eingebettet)."""
        slide = self.prs.slides.add_slide(self._layout())
        self._hintergrund_setzen(slide, self.farben["hintergrund"])

        pfad = Path(bild_pfad)
        if not pfad.exists():
            _err(f"Bilddatei nicht gefunden: {bild_pfad}")

        # SVG → PNG konvertieren
        if pfad.suffix.lower() == ".svg":
            png_bytes = _svg_zu_png(str(pfad))
            if png_bytes:
                import tempfile
                with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                    tmp.write(png_bytes)
                    bild_pfad = tmp.name
            else:
                _log(f"SVG-Konvertierung fehlgeschlagen — Folie wird ohne Bild erzeugt")
                self._textbox(slide, f"[Diagramm: {pfad.name}]",
                               x=0.3, y=1.5, breite=9.4, hoehe=2.5,
                               schriftgroesse=24, farbe_hex=self.farben["text_muted"],
                               ausrichtung="mitte")
                self.slide_count += 1
                return

        self._rechteck(slide, 0, 0, 10, 0.04, self.farben["akzent"])

        if titel:
            self._textbox(slide, titel,
                           x=0.3, y=0.1, breite=9.4, hoehe=0.5,
                           schriftgroesse=16, fett=True,
                           farbe_hex=self.farben["text"])

        bild_y = 0.65 if titel else 0.1
        bild_h = 5.625 - bild_y - (0.4 if beschriftung else 0.1)

        if vollflächig:
            slide.shapes.add_picture(bild_pfad,
                                      self._Inches(0.1), self._Inches(bild_y),
                                      self._Inches(9.8), self._Inches(bild_h))
        else:
            slide.shapes.add_picture(bild_pfad,
                                      self._Inches(0.3), self._Inches(bild_y),
                                      self._Inches(7), self._Inches(bild_h))
            if beschriftung:
                self._textbox(slide, beschriftung,
                               x=7.5, y=bild_y, breite=2.2, hoehe=bild_h,
                               schriftgroesse=10,
                               farbe_hex=self.farben["text_muted"])

        if beschriftung and vollflächig:
            self._textbox(slide, beschriftung,
                           x=0.3, y=5.1, breite=9.4, hoehe=0.4,
                           schriftgroesse=10, kursiv=True,
                           farbe_hex=self.farben["text_muted"],
                           ausrichtung="mitte")

        self.slide_count += 1
        _ok(f"Folie {self.slide_count}: Diagramm '{pfad.name}'")

    def slide_text(self, inhalt: str, titel: str = "", stichpunkte: bool = True):
        """
        Textfolie mit optionalen Stichpunkten.
        Markdown-Stichpunkte (- Text oder * Text) werden als Bullet-Points gerendert.
        """
        slide = self.prs.slides.add_slide(self._layout())
        self._hintergrund_setzen(slide, self.farben["hintergrund"])
        self._rechteck(slide, 0, 0, 10, 0.04, self.farben["akzent"])

        if titel:
            self._textbox(slide, titel,
                           x=0.4, y=0.1, breite=9.2, hoehe=0.55,
                           schriftgroesse=20, fett=True,
                           farbe_hex=self.farben["text"])

        zeilen = inhalt.strip().split("\n")
        y_start = 0.8 if titel else 0.3
        y_schritt = 0.42
        y = y_start

        for zeile in zeilen:
            zeile = zeile.strip()
            if not zeile:
                y += 0.2
                continue

            ist_bullet = zeile.startswith(("- ", "* ", "• "))
            text = zeile[2:] if ist_bullet else zeile
            einzug = 0.7 if ist_bullet else 0.4

            if ist_bullet:
                # Bullet-Punkt
                self._textbox(slide, "•",
                               x=0.4, y=y, breite=0.3, hoehe=y_schritt,
                               schriftgroesse=14,
                               farbe_hex=self.farben["akzent"])

            self._textbox(slide, text,
                           x=einzug, y=y, breite=9.2-einzug, hoehe=y_schritt,
                           schriftgroesse=14,
                           farbe_hex=self.farben["text"])
            y += y_schritt

            if y > 5.0:
                _log(f"Text-Folie: Inhalt überschreitet Folienhöhe — Rest abgeschnitten")
                break

        self.slide_count += 1
        _ok(f"Folie {self.slide_count}: Textfolie '{titel}' ({len(zeilen)} Zeilen)")

    def speichern(self, pfad: str) -> None:
        self.prs.save(pfad)
        _ok(f"Präsentation gespeichert: {pfad} ({self.slide_count} Folien)")

# =============================================================================
# MODUS 1: ERSTELLEN — Vollständige Präsentation aus mehreren Quellen
# =============================================================================

def modus_erstellen(titel: str, args: list[str]) -> None:
    """Erstellt vollständige Präsentation aus Optionen."""
    def _opt(name: str, default: str | None = None) -> str | None:
        for i, a in enumerate(args):
            if a == name and i + 1 < len(args):
                return args[i + 1]
        return default

    output     = _opt("--output",   f"{titel.replace(' ', '_')}.pptx")
    thema      = _opt("--thema",    DEFAULT_THEMA)
    template   = _opt("--template", None)
    untertitel = _opt("--untertitel", "")
    excel_pfad = _opt("--excel",    None)
    sheet      = _opt("--sheet",    None)
    diagramm   = _opt("--diagramm", None)
    kpi_json   = _opt("--kpis",     None)

    prs = Praesentation(thema=thema, template=template)

    # 1. Titelfolie
    prs.slide_titel(
        titel=titel,
        untertitel=untertitel,
        datum=datetime.now().strftime("%d. %B %Y")
    )

    # 2. KPI-Folie (aus JSON-Datei oder Inline-JSON)
    if kpi_json:
        try:
            if Path(kpi_json).exists():
                kpis = json.loads(Path(kpi_json).read_text())
            else:
                kpis = json.loads(kpi_json)
            prs.slide_abschnitt("Kennzahlen", "01")
            prs.slide_kpi(kpis)
        except Exception as e:
            _log(f"KPI-Parsing fehlgeschlagen: {e}")

    # 3. Excel-Daten (mehrere Sheets)
    if excel_pfad and Path(excel_pfad).exists():
        try:
            import openpyxl
            wb = openpyxl.load_workbook(excel_pfad, read_only=True, data_only=True)
            sheets = wb.sheetnames[:3]  # max. 3 Sheets einbinden
            wb.close()

            prs.slide_abschnitt("Datenanalyse", "02")

            for sheet_name in sheets:
                header, zeilen = _excel_sheet_lesen(excel_pfad, sheet_name)
                if header and zeilen:
                    prs.slide_tabelle(
                        header=header,
                        zeilen=zeilen,
                        titel=f"{sheet_name} — Auszug (max. 15 Zeilen)"
                    )
        except Exception as e:
            _log(f"Excel-Import fehlgeschlagen: {e}")

    # 4. Diagramm
    if diagramm and Path(diagramm).exists():
        prs.slide_abschnitt("Visualisierung", "03")
        prs.slide_diagramm(
            bild_pfad=diagramm,
            titel=Path(diagramm).stem.replace("_", " ").title(),
            vollflächig=True
        )

    # 5. Speichern
    prs.speichern(output)
    _json_out({
        "status":   "ok",
        "output":   output,
        "folien":   prs.slide_count,
        "thema":    thema,
        "quellen": {
            "excel":    excel_pfad,
            "diagramm": diagramm,
            "kpis":     kpi_json,
        },
    })

# =============================================================================
# MODUS 2: HINZUFÜGEN — Einzelne Folie zu bestehender Präsentation
# =============================================================================

def modus_hinzufuegen(pptx_pfad: str, slide_typ: str, args: list[str]) -> None:
    """Fügt eine einzelne Folie zu einer bestehenden Präsentation hinzu."""
    def _opt(name: str, default: str | None = None) -> str | None:
        for i, a in enumerate(args):
            if a == name and i + 1 < len(args):
                return args[i + 1]
        return default

    if not Path(pptx_pfad).exists():
        _err(f"Präsentation nicht gefunden: {pptx_pfad}")

    thema  = _opt("--thema", DEFAULT_THEMA)
    titel  = _opt("--titel", "")
    output = _opt("--output", pptx_pfad)  # Standard: in-place

    prs = Praesentation(thema=thema, template=pptx_pfad)

    if slide_typ == "kpi":
        kpi_json = _opt("--kpis", "[]")
        try:
            kpis = json.loads(kpi_json)
        except Exception:
            _err("--kpis muss valides JSON sein: '[{\"label\":\"...\",\"wert\":\"...\"}]'")
        prs.slide_kpi(kpis)

    elif slide_typ == "tabelle":
        csv_pfad = _opt("--csv", None)
        excel_pfad = _opt("--excel", None)
        if csv_pfad:
            header, zeilen = _csv_lesen(csv_pfad)
        elif excel_pfad:
            header, zeilen = _excel_sheet_lesen(excel_pfad, _opt("--sheet"))
        else:
            _err("--csv PFAD oder --excel PFAD erforderlich für Tabellen-Folie")
        prs.slide_tabelle(header=header, zeilen=zeilen, titel=titel)

    elif slide_typ == "diagramm":
        bild = _opt("--bild", None)
        if not bild:
            _err("--bild PFAD erforderlich für Diagramm-Folie")
        prs.slide_diagramm(bild_pfad=bild, titel=titel)

    elif slide_typ == "text":
        inhalt = _opt("--inhalt", "")
        if not inhalt:
            _err("--inhalt TEXT erforderlich für Textfolie")
        prs.slide_text(inhalt=inhalt, titel=titel)

    elif slide_typ == "abschnitt":
        prs.slide_abschnitt(kapitel=titel or "Abschnitt", nummer=_opt("--nummer", ""))

    else:
        _err(f"Unbekannter Slide-Typ: {slide_typ}\nErlaubt: kpi | tabelle | diagramm | text | abschnitt")

    prs.speichern(output)
    _json_out({"status": "ok", "output": output, "slide_typ": slide_typ})

# =============================================================================
# MODUS 3: VORSCHAU — Metadaten einer bestehenden Präsentation ausgeben
# =============================================================================

def modus_vorschau(pptx_pfad: str) -> None:
    """Gibt Metadaten und Folien-Übersicht einer bestehenden Präsentation aus."""
    if not Path(pptx_pfad).exists():
        _err(f"Präsentation nicht gefunden: {pptx_pfad}")

    Presentation, Inches, Pt, Emu, RGBColor, PP_ALIGN = _pptx_import()
    prs = Presentation(pptx_pfad)

    folien = []
    for i, slide in enumerate(prs.slides):
        texte = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texte.append(shape.text.strip()[:80])
        folien.append({
            "nummer":   i + 1,
            "layout":   slide.slide_layout.name,
            "texte":    texte[:5],
            "shapes":   len(slide.shapes),
        })

    _json_out({
        "datei":          pptx_pfad,
        "folienanzahl":   len(prs.slides),
        "breite_zoll":    round(prs.slide_width.inches, 2),
        "hoehe_zoll":     round(prs.slide_height.inches, 2),
        "seitenverhältnis": "16:9" if abs(prs.slide_width.inches / prs.slide_height.inches - 16/9) < 0.1 else "4:3",
        "folien":         folien,
    })

# =============================================================================
# HILFE
# =============================================================================

HILFE = """
VERWENDUNG:
  python pptx_generator.py erstellen TITEL [OPTIONEN]
  python pptx_generator.py hinzufuegen PFAD SLIDE-TYP [OPTIONEN]
  python pptx_generator.py vorschau PFAD

ERSTELLEN — Optionen:
  --output PFAD         Ausgabedatei (Standard: TITEL.pptx)
  --thema dunkel|hell|corporate  Farbschema (Standard: dunkel)
  --template PFAD       Corporate-Vorlage (.pptx als Basis)
  --untertitel TEXT     Untertitel auf Titelfolie
  --excel PFAD          Excel-Workbook einbinden (Sheets → Tabellen-Folien)
  --sheet NAME          Nur dieses Sheet aus Excel verwenden
  --diagramm PFAD       SVG oder PNG auf Diagramm-Folie
  --kpis JSON           KPI-Array: '[{"label":"Umsatz","wert":"€2.4M","trend":"+12%","positiv":true}]'
                        oder Pfad zu JSON-Datei

HINZUFÜGEN — Slide-Typen:
  kpi       --kpis JSON
  tabelle   --csv PFAD | --excel PFAD [--sheet NAME] [--titel TEXT]
  diagramm  --bild PFAD [--titel TEXT]
  text      --inhalt "TEXT" [--titel TEXT]
  abschnitt --titel TEXT [--nummer 01]

VORAUSSETZUNG:
  pip install python-pptx
  pip install cairosvg pillow      (für SVG→PNG-Konvertierung)
  pip install openpyxl             (für Excel-Import)

HYBRID-WORKFLOW-BEISPIEL (Quartalsbericht):
  # Schritt 3a + 3b parallel:
  python excel_finanzmodell.py q1_bereinigt.csv --output q1_bericht.xlsx &
  python diagramm_generator.py sankey q1_bereinigt.csv --output fluss.svg &
  wait

  # Schritt 4: PowerPoint aus beiden Quellen
  python pptx_generator.py erstellen "Q1 2026 Quartalsbericht" \\
    --excel    q1_bericht.xlsx \\
    --diagramm fluss.svg \\
    --kpis     '[{"label":"Gesamtumsatz","wert":"€ 4.2M","trend":"+8%","positiv":true}, \\
                 {"label":"Transaktionen","wert":"12.847","trend":"+234","positiv":true}, \\
                 {"label":"Fehlerrate","wert":"0.3%","trend":"-0.1pp","positiv":true}]' \\
    --thema    corporate \\
    --output   Q1_2026_Quartalsbericht.pptx
"""

# =============================================================================
# CLI
# =============================================================================

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print(HILFE); sys.exit(0)

    modus = sys.argv[1]
    rest  = sys.argv[2:]

    if modus == "erstellen":
        if not rest or rest[0].startswith("--"):
            _err("VERWENDUNG: erstellen TITEL [OPTIONEN]")
        modus_erstellen(rest[0], rest[1:])

    elif modus == "hinzufuegen":
        if len(rest) < 2:
            _err("VERWENDUNG: hinzufuegen PFAD SLIDE-TYP [OPTIONEN]")
        modus_hinzufuegen(rest[0], rest[1], rest[2:])

    elif modus == "vorschau":
        if not rest:
            _err("VERWENDUNG: vorschau PFAD")
        modus_vorschau(rest[0])

    else:
        print(f"ERR | Unbekannter Modus: '{modus}'\n{HILFE}", file=sys.stderr)
        sys.exit(1)
